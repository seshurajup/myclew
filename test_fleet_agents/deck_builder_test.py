"""deck_builder_test — data-wise verifier for the ppt-master-distilled spec→.pptx builder.

Core properties:
  1. build_deck writes a real .pptx with the expected slide count (title + N content slides), read back.
  2. bullets and a table land in the deck (text present).
  3. ledger_to_spec turns experiment rows into a valid deck spec.
  4. agent contract."""
import os, sys, tempfile
COMP = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, COMP); sys.path.insert(0, os.path.join(COMP, "tools", "researchpapers"))
from fleet_agents import deck_builder as D


def _run():
    print("=== DECK-BUILDER VERIFIER ===")
    checks = {}
    try:
        from pptx import Presentation
    except Exception as e:
        print("python-pptx missing:", e); return False

    out = os.path.join(tempfile.mkdtemp(), "t.pptx")
    spec = {"title": "T", "subtitle": "S",
            "slides": [{"heading": "H1", "bullets": ["b1", "b2"]},
                       {"heading": "H2", "bullets": ["x"], "table": [["m", "v"], ["cv", "0.88"]]}]}
    r = D.build_deck(spec, out)
    checks["file_written"] = os.path.exists(out) and os.path.getsize(out) > 0
    checks["slide_count"] = r["n_slides"] == 3            # 1 title + 2 content

    prs = Presentation(out)
    alltext = " ".join(sh.text_frame.text for sl in prs.slides for sh in sl.shapes if sh.has_text_frame)
    checks["bullets_present"] = "b1" in alltext and "b2" in alltext
    tabletext = " ".join(c.text for sl in prs.slides for sh in sl.shapes if sh.has_table
                         for row in sh.table.rows for c in row.cells)
    checks["table_present"] = "0.88" in tabletext
    print(f"  -> wrote {r['n_slides']} slides; bullets+table read back OK")

    # ledger_to_spec
    sp = D.ledger_to_spec([{"change": "muon", "cv": 0.88, "lb": 0.87, "description": "d"}], "Ledger")
    checks["ledger_spec"] = sp["slides"][0]["table"][0] == ["change", "cv", "lb"] and len(sp["slides"]) == 2

    # agent
    st, dta, to, msg = D.run_deck({"spec": {}}, "t")
    checks["agent_done"] = st == "done" and os.path.exists(dta["path"])

    for k, v in checks.items():
        print(f"  {'OK' if v else 'X'} {k}")
    ok = all(checks.values())
    print(f"=== deck-builder: {'PASS' if ok else 'FAIL'} ({sum(checks.values())}/{len(checks)}) ===")
    return ok


if __name__ == "__main__":
    sys.exit(0 if _run() else 1)
