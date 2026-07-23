"""deck_builder — structured-spec → .pptx report generator, the offline/deterministic primitive distilled
from hugohe3/ppt-master. ppt-master itself is an LLM agent that hand-authors SVG per slide then compiles to
DrawingML — powerful but heavy and non-deterministic. For a fleet the useful, reusable piece is a thin
spec→deck helper on python-pptx (already this env's dependency): give it a title + a list of slides (heading,
bullets, optional table), get a real .pptx — deterministic, no LLM, no network. Used to turn an experiment
ledger / CV summary / writeup into a shareable deck.

Spec shape:
  {"title": str, "subtitle": str,
   "slides": [ {"heading": str, "bullets": [str, ...], "table": [[...rows...]] (optional)} , ... ]}

Primitives (deps = python-pptx, import-guarded):
  • build_deck(spec, out_path)   — write the .pptx; returns the path and slide count.
  • ledger_to_spec(rows, title)  — convenience: turn experiment-ledger dicts into a deck spec.
"""
from __future__ import annotations
from .base import BaseAgent


def build_deck(spec, out_path):
    """Render `spec` to a .pptx at out_path. Returns {"path", "n_slides"}. Requires python-pptx."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    # title slide
    ts = prs.slides.add_slide(prs.slide_layouts[0])
    ts.shapes.title.text = str(spec.get("title", "Report"))
    if spec.get("subtitle") and len(ts.placeholders) > 1:
        ts.placeholders[1].text = str(spec["subtitle"])

    for sl in spec.get("slides", []):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = str(sl.get("heading", ""))
        body = slide.placeholders[1].text_frame
        body.clear()
        bullets = sl.get("bullets", []) or []
        for i, b in enumerate(bullets):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = str(b); p.level = 0
        table = sl.get("table")
        if table:
            rows, cols = len(table), max(len(r) for r in table)
            gt = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(4.2),
                                        Inches(9), Inches(0.4 * rows)).table
            for r, row in enumerate(table):
                for c in range(cols):
                    gt.cell(r, c).text = str(row[c]) if c < len(row) else ""
                    gt.cell(r, c).text_frame.paragraphs[0].font.size = Pt(12)
    prs.save(str(out_path))
    return {"path": str(out_path), "n_slides": len(prs.slides._sldIdLst)}


def ledger_to_spec(rows, title="Experiment Ledger"):
    """Turn a list of experiment dicts (change/cv/lb/description) into a deck spec: a summary table slide +
    one detail slide per row."""
    header = ["change", "cv", "lb"]
    table = [header] + [[str(r.get(k, "")) for k in header] for r in rows]
    slides = [{"heading": "Summary", "bullets": [f"{len(rows)} experiments"], "table": table}]
    for r in rows:
        slides.append({"heading": str(r.get("change", "exp")),
                       "bullets": [f"CV: {r.get('cv','?')}", f"LB: {r.get('lb','?')}",
                                   str(r.get("description", ""))]})
    return {"title": title, "subtitle": f"{len(rows)} experiments", "slides": slides}


# ---------------------------------------------------------------- agent
class DeckBuilder(BaseAgent):
    name = "deck-builder"
    thread = "S"; kind = "finding"

    def run(self, q, worker):
        s = self.spec(q)
        try:
            import pptx  # noqa: F401
        except Exception as e:  # noqa: BLE001
            return self.escalate(q, "leader", f"deck-builder needs python-pptx: {e}")
        from .base import AUTO
        AUTO.mkdir(parents=True, exist_ok=True)
        out = s.get("out") or str(AUTO / "deck_builder_demo.pptx")
        spec = s.get("deck") or {"title": "Fleet Report", "subtitle": "auto",
                                 "slides": [{"heading": "Findings", "bullets": ["a", "b", "c"],
                                             "table": [["metric", "val"], ["cv", "0.88"]]}]}
        r = build_deck(spec, out)
        msg = (f"deck-builder: wrote {r['n_slides']}-slide .pptx → {r['path']}. Deterministic spec→deck "
               f"(title/bullets/table) on python-pptx — turn a ledger/CV summary into a shareable report "
               f"(ppt-master distilled, no LLM/SVG)")
        self.log(msg, kind="finding", recommendation="use ledger_to_spec(rows) then build_deck for writeup decks")
        return self.done({"path": r["path"], "n_slides": r["n_slides"]}, msg)


_AGENT = DeckBuilder()


def run_deck(q, worker):
    return _AGENT.run(q, worker)
